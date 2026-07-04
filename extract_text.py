"""
Извлечение текста из PDF/DOCX/DOCM/PPTX постранично (или послайдово), с
детектором "битого" текстового слоя (cid-коды вместо текста, как у части
страниц журнала «Горная промышленность») и реальным OCR-фолбэком через
Tesseract для таких страниц.

pip install pdfplumber python-docx python-pptx pytesseract PyMuPDF --break-system-packages
Плюс системный Tesseract с русским языковым пакетом:
    brew install tesseract tesseract-lang

Для .doc (старый бинарный формат Word 97-2003, НЕ OOXML — python-docx его
принципиально не читает) и для "битых"/нестандартных .docx/.docm нужен
LibreOffice как фолбэк-конвертер:
    brew install --cask libreoffice
Без LibreOffice такие файлы будут пропущены с понятным сообщением об ошибке.
"""
from __future__ import annotations

import logging
import re
import shutil
import subprocess
import tempfile
import threading
from dataclasses import dataclass, field
from pathlib import Path

import config
import yandex_ocr
from budget import BudgetExceededError

logger = logging.getLogger("extract_text")

# cid-коды выглядят как "(cid:123)" в сыром выводе некоторых экстракторов,
# либо после нормализации текст превращается в мусорные управляющие символы.
_CID_RE = re.compile(r"\(cid:\d+\)")

# Если Tesseract не установлен, не пытаемся OCR-ить каждую страницу заново
# (это дорого и бессмысленно) — фиксируем один раз и молча пропускаем OCR
# до конца прогона, оставляя needs_ocr=True как есть.
_tesseract_available: bool | None = None

# Бюджет на платный Yandex Vision OCR — простой счётчик на уровне модуля
# (не завязан на CLI, читает config.MAX_OCR_BUDGET_RUB). При исчерпании
# пайплайн НЕ останавливается, просто перестаёт звать Vision OCR и
# работает дальше только через бесплатный Tesseract.
_yandex_ocr_spent_rub = 0.0
_yandex_ocr_budget_exhausted = False
# При NUM_WORKERS > 1 несколько потоков могут одновременно читать/менять
# _yandex_ocr_spent_rub — без лока это гонка данных (можно проскочить лимит
# или потерять часть учтённых трат).
_yandex_ocr_lock = threading.Lock()


def _yandex_ocr_check_and_record() -> None:
    """Поднимает BudgetExceededError, если следующий вызов превысит MAX_OCR_BUDGET_RUB."""
    global _yandex_ocr_spent_rub
    with _yandex_ocr_lock:
        if _yandex_ocr_spent_rub + config.YANDEX_OCR_PRICE_PER_PAGE_RUB > config.MAX_OCR_BUDGET_RUB:
            raise BudgetExceededError(
                f"OCR-бюджет исчерпан: потрачено {_yandex_ocr_spent_rub:.2f}₽ "
                f"из {config.MAX_OCR_BUDGET_RUB:.2f}₽"
            )
        _yandex_ocr_spent_rub += config.YANDEX_OCR_PRICE_PER_PAGE_RUB


@dataclass
class PageText:
    page_number: int
    text: str
    needs_ocr: bool
    ocr_used: bool = False


@dataclass
class ExtractedText:
    full_text: str
    pages: list[PageText] = field(default_factory=list)


def _broken_ratio(text: str) -> float:
    if not text:
        return 1.0
    cid_hits = len(_CID_RE.findall(text))
    non_printable = sum(1 for ch in text if ord(ch) < 32 and ch not in "\n\t")
    junk_units = cid_hits + non_printable
    return min(1.0, junk_units / max(1, len(text) / 20))


def _tesseract_ready() -> bool:
    global _tesseract_available
    if _tesseract_available is not None:
        return _tesseract_available
    _tesseract_available = shutil.which("tesseract") is not None
    if not _tesseract_available:
        logger.warning(
            "Tesseract не найден в PATH — OCR-фолбэк для битых страниц PDF "
            "отключён на весь прогон. Установите: brew install tesseract tesseract-lang"
        )
    return _tesseract_available


def _png_to_jpeg_for_upload(png_bytes: bytes, max_dim: int = 2200, quality: int = 85) -> bytes:
    """
    Для Yandex Vision OCR не нужен полноразмерный 300dpi PNG — это лишние
    мегабайты, из-за которых загрузка не успевает уложиться в таймаут на
    больших/высоких страницах. Сжимаем в JPEG и слегка уменьшаем, оставляя
    разрешение, достаточное для распознавания текста.
    """
    from PIL import Image
    import io

    img = Image.open(io.BytesIO(png_bytes)).convert("RGB")
    if max(img.size) > max_dim:
        ratio = max_dim / max(img.size)
        img = img.resize((int(img.width * ratio), int(img.height * ratio)))
    out = io.BytesIO()
    img.save(out, format="JPEG", quality=quality)
    return out.getvalue()


def ocr_page_fallback(pdf_path: Path, page_number: int) -> str:
    """
    Рендерит страницу PDF в изображение (PyMuPDF, 300 dpi) и распознаёт
    текст. Основной движок — Yandex Vision OCR (лучше качество, сохраняет
    структуру таблиц в Markdown, не требует локальной установки); при
    ошибке сети/API или исчерпании OCR-бюджета — бесплатный откат на
    локальный Tesseract. Если недоступно вообще ничего — возвращает "".
    """
    global _yandex_ocr_budget_exhausted

    import fitz  # PyMuPDF

    try:
        doc = fitz.open(str(pdf_path))
        page = doc[page_number - 1]
        pix = page.get_pixmap(dpi=300)
        img_bytes = pix.tobytes("png")
        doc.close()
    except Exception as e:
        logger.warning("Не удалось отрендерить страницу %s стр. %d: %s", pdf_path.name, page_number, e)
        return ""

    # 1) Yandex Vision OCR — основной движок. Загружаем сжатый JPEG, а не
    # полноразмерный PNG, чтобы не упираться в таймаут на больших страницах.
    if config.YANDEX_API_KEY and not _yandex_ocr_budget_exhausted:
        try:
            _yandex_ocr_check_and_record()
            jpeg_bytes = _png_to_jpeg_for_upload(img_bytes)
            text = yandex_ocr.recognize_image(jpeg_bytes, mime_type="image/jpeg")
            if text.strip():
                return text
        except BudgetExceededError as e:
            logger.warning("%s Дальше использую бесплатный Tesseract.", e)
            _yandex_ocr_budget_exhausted = True
        except Exception as e:
            logger.warning(
                "Yandex Vision OCR не удался для %s стр. %d: %s, пробую Tesseract",
                pdf_path.name, page_number, e,
            )

    # 2) Tesseract — бесплатный фолбэк
    if not _tesseract_ready():
        return ""
    try:
        import pytesseract
        from PIL import Image
        import io

        img = Image.open(io.BytesIO(img_bytes))
        return pytesseract.image_to_string(img, lang="rus+eng")
    except Exception as e:
        logger.warning("Tesseract тоже не справился для %s стр. %d: %s", pdf_path.name, page_number, e)
        return ""


def extract_from_pdf(path: Path) -> ExtractedText:
    import pdfplumber

    pages: list[PageText] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            raw = page.extract_text() or ""
            ratio = _broken_ratio(raw)
            needs_ocr = ratio > config.BROKEN_TEXT_LAYER_THRESHOLD
            pages.append(PageText(page_number=i, text=raw, needs_ocr=needs_ocr))

    # OCR-проход по страницам с битым текстовым слоем — отдельно от
    # основного цикла, чтобы pdfplumber успел закрыть файл (OCR открывает
    # его заново через PyMuPDF).
    ocr_count = 0
    for p in pages:
        if not p.needs_ocr:
            continue
        ocr_text = ocr_page_fallback(path, p.page_number)
        if ocr_text.strip():
            p.text = ocr_text
            p.ocr_used = True
            p.needs_ocr = False
            ocr_count += 1
    if ocr_count:
        logger.info("OCR применён к %d стр. из %s", ocr_count, path.name)

    full_text = "\n\n".join(p.text for p in pages)
    return ExtractedText(full_text=full_text, pages=pages)


def _find_soffice() -> str | None:
    """
    На macOS LibreOffice.app не добавляет `soffice` в PATH автоматически —
    бинарник лежит внутри бандла приложения. Проверяем PATH, а если не
    нашли — известные стандартные пути установки (macOS/Linux).
    """
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    for candidate in (
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",  # macOS
        "/usr/bin/soffice",                                        # Linux (apt)
        "/usr/bin/libreoffice",
        "/opt/libreoffice/program/soffice",
    ):
        if Path(candidate).exists():
            return candidate
    return None


def _convert_via_libreoffice(path: Path) -> Path | None:
    """
    Фолбэк для файлов, которые python-docx не может открыть напрямую:
    legacy .doc (бинарный формат Word 97-2003, не OOXML), .docm с
    нестандартным content-type, повреждённые/нестандартные .doc(x).
    Конвертирует через headless LibreOffice в .docx во временную папку.

    Возвращает путь к сконвертированному .docx или None, если LibreOffice
    не найден / конвертация не удалась.
    """
    soffice = _find_soffice()
    if not soffice:
        return None
    tmp_dir = Path(tempfile.mkdtemp(prefix="loconv_"))
    # Каждый вызов — свой изолированный профиль LibreOffice: без этого
    # параллельные конвертации (NUM_WORKERS > 1) конфликтуют на общем
    # профиле и падают/зависают.
    profile_dir = Path(tempfile.mkdtemp(prefix="lo_profile_"))
    try:
        subprocess.run(
            [
                soffice, "--headless",
                f"-env:UserInstallation=file://{profile_dir}",
                "--convert-to", "docx", "--outdir", str(tmp_dir), str(path),
            ],
            check=True, capture_output=True, timeout=120,
        )
        converted = tmp_dir / (path.stem + ".docx")
        return converted if converted.exists() else None
    except Exception:
        return None


def resolve_docx_path(path: Path) -> Path:
    """
    Возвращает путь к файлу, который python-docx точно сможет открыть:
    - если исходный файл уже открывается нормально — возвращает его же;
    - если нет (старый бинарный .doc, .docm с macro-content-type, битый
      .docx) — конвертирует через LibreOffice и возвращает путь к
      сконвертированному .docx.
    Вызывать ОДИН раз на файл и передавать результат во все экстракторы
    (текст/таблицы/картинки), а не звать фолбэк в каждом из них отдельно —
    иначе LibreOffice будет молотить один и тот же файл по три раза.
    """
    import docx

    try:
        docx.Document(str(path))
        return path
    except Exception as e:
        converted = _convert_via_libreoffice(path)
        if converted is None:
            raise ValueError(
                f"Не удалось прочитать {path.name} напрямую ({e}). Похоже, "
                f"это старый бинарный .doc (не OOXML) или повреждённый файл, "
                f"а LibreOffice для фолбэка не найден/не смог сконвертировать."
            ) from e
        return converted


def extract_from_docx(path: Path) -> ExtractedText:
    import docx

    resolved = resolve_docx_path(path)
    d = docx.Document(str(resolved))

    paragraphs = [p.text for p in d.paragraphs if p.text.strip()]
    full_text = "\n".join(paragraphs)
    return ExtractedText(
        full_text=full_text,
        pages=[PageText(page_number=1, text=full_text, needs_ocr=False)],
    )


def extract_from_pptx(path: Path) -> ExtractedText:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[PageText] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        chunks.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.append(" | ".join(cell.text for cell in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                chunks.append(f"[Заметки] {notes}")
        pages.append(PageText(page_number=i, text="\n".join(chunks), needs_ocr=False))

    full_text = "\n\n".join(f"[Слайд {p.page_number}]\n{p.text}" for p in pages if p.text.strip())
    return ExtractedText(full_text=full_text, pages=pages)


def extract_text(path: Path) -> ExtractedText:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_from_pdf(path)
    if suffix in (".docx", ".doc", ".docm"):
        return extract_from_docx(path)
    if suffix == ".pptx":
        return extract_from_pptx(path)
    raise ValueError(f"Неподдерживаемый формат для извлечения текста: {suffix}")
